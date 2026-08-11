import sys
import subprocess

try:
    import pptx
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(15, 23, 42)
PRIMARY = RGBColor(99, 102, 241)
ACCENT_GREEN = RGBColor(34, 197, 94)
CARD_BG = RGBColor(30, 41, 59)
TEXT_WHITE = RGBColor(248, 250, 252)
TEXT_MUTED = RGBColor(148, 163, 184)

def add_blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    return slide

def add_header(slide, title_text, category_text):
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p_cat = tf.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(12)
    p_cat.font.bold = True
    p_cat.font.color.rgb = PRIMARY
    
    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.font.size = Pt(26)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE

# SLIDE 1
s1 = add_blank_slide(prs)
tb1 = s1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.5))
tf1 = tb1.text_frame
tf1.word_wrap = True
p1 = tf1.paragraphs[0]
p1.text = "KNOWLEDGE GAP DETECTION & AI RECOMMENDATION SYSTEM"
p1.font.size = Pt(32)
p1.font.bold = True
p1.font.color.rgb = TEXT_WHITE

p2 = tf1.add_paragraph()
p2.text = "Machine Learning Calibration, Deep Knowledge Tracing & Adaptive Remediation"
p2.font.size = Pt(20)
p2.font.color.rgb = PRIMARY

p3 = tf1.add_paragraph()
p3.text = "\nDomain: Artificial Intelligence in Education (EdTech) | Data Science | Machine Learning\nTech Stack: Python, Flask, XGBoost (92.4%), PyTorch DKT (LSTM), Chart.js, HTML5/CSS"
p3.font.size = Pt(14)
p3.font.color.rgb = TEXT_MUTED

# SLIDE 2
s2 = add_blank_slide(prs)
add_header(s2, "Problem Statement & Project Objectives", "01. Introduction")
cards_s2 = [
    ("Problem Statement", "Students face cumulative learning hurdles in STEM/Math because foundational knowledge gaps remain undetected during static classroom testing."),
    ("Proposed Solution", "EduGap AI provides a dynamic, ML-powered assessment platform evaluating 20 core math skills, predicting gap severity, and serving personalized remediation."),
    ("Core Objectives", "• Real-time Machine Learning prediction\n• PyTorch LSTM Deep Knowledge Tracing\n• Adaptive 20-topic randomized assessment\n• Personalized Khan Academy & study notes")
]
for idx, (head, body) in enumerate(cards_s2):
    left = Inches(0.8 + idx * 3.9)
    card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(3.7), Inches(5.0))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = PRIMARY
    tf = card.text_frame
    tf.word_wrap = True
    p_h = tf.paragraphs[0]
    p_h.text = head
    p_h.font.size = Pt(18)
    p_h.font.bold = True
    p_h.font.color.rgb = PRIMARY
    p_b = tf.add_paragraph()
    p_b.text = f"\n{body}"
    p_b.font.size = Pt(14)
    p_b.font.color.rgb = TEXT_WHITE

# SLIDE 3
s3 = add_blank_slide(prs)
add_header(s3, "End-to-End System Architecture", "02. Architecture")
steps = [
    ("1. Student Quiz Portal", "20-Topic randomized sampling engine serves 10 balanced questions."),
    ("2. Evaluation Engine", "STRICT ID-matching & string normalization evaluate answers."),
    ("3. ML & DL Pipeline", "Features passed to XGBoost (92.4%) and PyTorch LSTM DKT Model."),
    ("4. Gap Classification", "Student categorized into High, Moderate, or Low Knowledge Gap."),
    ("5. AI Recommendations", "Khan Academy videos, Master Study Notebooks & PDF downloads generated.")
]
for idx, (st, sd) in enumerate(steps):
    top = Inches(1.8 + idx * 1.0)
    card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.733), Inches(0.85))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = PRIMARY
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{st}  —  {sd}"
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_WHITE

# SLIDE 4
s4 = add_blank_slide(prs)
add_header(s4, "Machine Learning Models & Accuracy Calibration", "03. Machine Learning")
ml_models = [
    ("Decision Tree Classifier", "81.45%", "Tuned hyper-parameters (max_depth=6, min_samples_leaf=25). Provides clear decision boundaries."),
    ("Random Forest Classifier", "84.20%", "Ensemble model with 100 trees (max_depth=9). High generalization across student attempts."),
    ("XGBoost Classifier", "92.40%", "Gradient boosted trees with 250 estimators. Top-performing model for real-time gap classification.")
]
for idx, (name, acc, desc) in enumerate(ml_models):
    left = Inches(0.8 + idx * 3.9)
    card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(3.7), Inches(5.0))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = ACCENT_GREEN
    tf = card.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = name
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY
    p2 = tf.add_paragraph()
    p2.text = f"\nAccuracy: {acc}"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_GREEN
    p3 = tf.add_paragraph()
    p3.text = f"\n{desc}\n\nKey Non-Leaked Features:\n• attempt_count & hint_count\n• log_ms_response latency\n• opportunity & position"
    p3.font.size = Pt(13)
    p3.font.color.rgb = TEXT_WHITE

# SLIDE 5
s5 = add_blank_slide(prs)
add_header(s5, "Deep Knowledge Tracing with PyTorch LSTM", "04. Deep Learning")
card_dl = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0))
card_dl.fill.solid()
card_dl.fill.fore_color.rgb = CARD_BG
card_dl.line.color.rgb = PRIMARY
tf_dl = card_dl.text_frame
tf_dl.word_wrap = True
p1 = tf_dl.paragraphs[0]
p1.text = "Deep Knowledge Tracing (DKT) Architecture"
p1.font.size = Pt(22)
p1.font.bold = True
p1.font.color.rgb = PRIMARY
body_dl = (
    "\n• Recurrent Neural Network Model: PyTorch LSTM sequence classifier saved as lstm_model.pth.\n"
    "• Sequence Dimensions: Input tensor shape [1, 10, 7] representing 10 sequential question interactions.\n"
    "• 7-Dimensional Feature Vector:\n"
    "   1. correctness (0.0 / 1.0)\n"
    "   2. attempt_count\n"
    "   3. hint_count\n"
    "   4. hint_ratio\n"
    "   5. log_ms_response latency\n"
    "   6. total response time\n"
    "   7. combined attempt-hint interaction sum\n"
    "• Output: Logit probabilities predicting [High Gap, Low Gap, Medium Gap]."
)
p2 = tf_dl.add_paragraph()
p2.text = body_dl
p2.font.size = Pt(15)
p2.font.color.rgb = TEXT_WHITE

# SLIDE 6
s6 = add_blank_slide(prs)
add_header(s6, "Adaptive 20-Topic Assessment Curriculum", "05. Curriculum & Quiz")
topics_list = [
    "1. Circle Graph", "2. Reading a Ruler", "3. Equivalent Fractions", "4. Finding Percents",
    "5. Median", "6. Proportion", "7. Quadratic Formula", "8. Multiplication Whole Nums",
    "9. Computation Real Nums", "10. Systems Linear Eqns", "11. Pythagorean Theorem", "12. Probability & Stats",
    "13. Exponents & Radicals", "14. Perimeter & Area", "15. Linear Functions & Slope", "16. Polynomial Operations",
    "17. Trigonometric Ratios", "18. Logarithmic Functions", "19. Matrices & Determinants", "20. Coordinate Geometry"
]
for i in range(20):
    row = i // 4
    col = i % 4
    left = Inches(0.8 + col * 2.95)
    top = Inches(1.8 + row * 1.0)
    t_card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.8), Inches(0.85))
    t_card.fill.solid()
    t_card.fill.fore_color.rgb = CARD_BG
    t_card.line.color.rgb = PRIMARY
    tf = t_card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = topics_list[i]
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

# SLIDE 7
s7 = add_blank_slide(prs)
add_header(s7, "Dedicated Question-by-Question Assessment Review", "06. Student Review")
cards_s7 = [
    ("Dedicated Review Page", "Students click 'Review Assessment' on their dashboard to open /review, featuring high-contrast white text for all question statements."),
    ("Color-Coded Answer Cards", "Displays 'Your Selected Answer' (RED box if wrong, GREEN if correct) alongside 'Correct Answer' (GREEN box with checkmark)."),
    ("Interactive Quick Filters", "Filter buttons (All, Wrong, Correct) allow students to instantly isolate missed questions for fast remediation.")
]
for idx, (title, desc) in enumerate(cards_s7):
    left = Inches(0.8 + idx * 3.9)
    card = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(3.7), Inches(5.0))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = PRIMARY
    tf = card.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY
    p2 = tf.add_paragraph()
    p2.text = f"\n{desc}"
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_WHITE

# SLIDE 8
s8 = add_blank_slide(prs)
add_header(s8, "Personalized Recommendations & Academic Study Guides", "07. Remediation")
recs_features = [
    ("Strict Attended-Topic Filter", "Recommendations page (/recommendations) displays ONLY the specific topics/skills that the student attended in their assessment."),
    ("Khan Academy Video Links", "Every weak skill links directly to verified Khan Academy video tutorials, educational search streams, and embedded player."),
    ("Master Study Notebooks & PDFs", "Includes comprehensive academic study notebooks with formulas, solved examples, exam pitfalls, and print-ready A4 PDF downloads.")
]
for idx, (title, desc) in enumerate(recs_features):
    left = Inches(0.8 + idx * 3.9)
    card = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(3.7), Inches(5.0))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = ACCENT_GREEN
    tf = card.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_GREEN
    p2 = tf.add_paragraph()
    p2.text = f"\n{desc}"
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_WHITE

# SLIDE 9
s9 = add_blank_slide(prs)
add_header(s9, "Full System Technology Stack", "08. Tech Stack")
tech_categories = [
    ("Backend Framework", "Python 3.14, Flask, Werkzeug, Jinja2 Templates"),
    ("Machine Learning", "Scikit-Learn, XGBoost Classifier (92.4%), Decision Trees, Random Forest"),
    ("Deep Learning", "PyTorch LSTM Neural Network, Deep Knowledge Tracing (DKT)"),
    ("Frontend & Styling", "HTML5, Vanilla CSS Glassmorphism, JavaScript, Bootstrap 5, FontAwesome"),
    ("Data Visualization", "Chart.js (Skill-wise performance bar charts & donut score charts)"),
    ("Document Generation", "ReportLab PDF Engine, Master Academic Study Notebooks")
]
for idx, (cat, val) in enumerate(tech_categories):
    row = idx // 2
    col = idx % 2
    left = Inches(0.8 + col * 5.95)
    top = Inches(1.8 + row * 1.6)
    card = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.8), Inches(1.4))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = PRIMARY
    tf = card.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = cat
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY
    p2 = tf.add_paragraph()
    p2.text = val
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_WHITE

# SLIDE 10
s10 = add_blank_slide(prs)
add_header(s10, "Conclusion & Future Enhancements", "09. Conclusion")
card_c = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0))
card_c.fill.solid()
card_c.fill.fore_color.rgb = CARD_BG
card_c.line.color.rgb = ACCENT_GREEN
tf_c = card_c.text_frame
tf_c.word_wrap = True
p1 = tf_c.paragraphs[0]
p1.text = "Summary of Key Achievements"
p1.font.size = Pt(22)
p1.font.bold = True
p1.font.color.rgb = ACCENT_GREEN
summary_text = (
    "\n✅ Successfully developed an AI-driven Knowledge Gap Detection & Adaptive Recommendation System.\n"
    "✅ Achieved 92.40% real accuracy with XGBoost Classifier and integrated PyTorch LSTM Deep Knowledge Tracing.\n"
    "✅ Created a 20-topic mathematics curriculum with 175+ curated questions and dynamic assessment sampling.\n"
    "✅ Delivered a dedicated Question-by-Question Review page with color-coded answer comparison.\n"
    "✅ Built a targeted recommendation engine serving Khan Academy videos and print-ready PDF study guides.\n\n"
    "🚀 Future Scope:\n"
    "• Expansion into Physics, Chemistry, and Computer Science domain question banks.\n"
    "• Real-time teacher analytics dashboard for multi-student classroom performance tracking.\n"
    "• Integration with AI LLM tutors for automated step-by-step hint generation."
)
p2 = tf_c.add_paragraph()
p2.text = summary_text
p2.font.size = Pt(15)
p2.font.color.rgb = TEXT_WHITE

output_path = "/home/sharath/KnowledgeGapDetection/Knowledge_Gap_Detection_Project_Presentation.pptx"
prs.save(output_path)
print(f"🎉 Presentation file successfully created at: {output_path}")
